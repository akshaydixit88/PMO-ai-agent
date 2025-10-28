from pptx import Presentation
from pptx.dml.color import RGBColor
from datetime import datetime
import re

def generate_dashboard_from_exec(exec_summary_text, template_path="template.pptx"):
    """
    Generates a dashboard slide from a raw exec summary text using a template.
    """
    # ---- Parse the exec summary ----
    # Status
    status_match = re.search(r"\*\*Overall Program Status:\*\*\s*(.+)", exec_summary_text)
    status = status_match.group(1).strip() if status_match else "Gray"

    # Progress
    progress_match = re.search(r"\*\*Key Progress Points:\*\*(.*?)\n\n\*\*Risks", exec_summary_text, re.S)
    progress = progress_match.group(1).strip() if progress_match else ""

    # Risks
    risks_match = re.search(r"\*\*Risks and Mitigation Actions:\*\*(.*)", exec_summary_text, re.S)
    risks = risks_match.group(1).strip() if risks_match else ""

    # Summary (short final takeaway)
    summary = "Program updates automatically generated."

    update_data = {
        "status": status,
        "progress": progress,
        "risks": risks,
        "summary": summary
    }

    # ---- Helper to get shape by name ----
    def get_shape_by_name(slide, name):
        for shape in slide.shapes:
            if shape.name == name:
                return shape
        raise ValueError(f"Shape with name '{name}' not found")

    # ---- Load template ----
    prs = Presentation(template_path)
    slide = prs.slides[0]

    # ---- Get shapes ----
    title_shape = get_shape_by_name(slide, 'title_box')
    status_shape = get_shape_by_name(slide, 'status_rect')
    progress_shape = get_shape_by_name(slide, 'progress_box')
    risks_shape = get_shape_by_name(slide, 'risks_box')
    summary_shape = get_shape_by_name(slide, 'summary_box')

    # ---- Populate title ----
    today_str = datetime.now().strftime("%b %d, %Y")
    title_shape.text = f"Program Update - {today_str}"

    # ---- Status rectangle color ----
    fill = status_shape.fill
    fill.solid()
    status_lower = status.lower()
    if status_lower == "green":
        fill.fore_color.rgb = RGBColor(0, 128, 0)
    elif status_lower == "yellow":
        fill.fore_color.rgb = RGBColor(255, 204, 0)
    elif status_lower == "red":
        fill.fore_color.rgb = RGBColor(255, 0, 0)
    else:
        fill.fore_color.rgb = RGBColor(128, 128, 128)

    # ---- Populate text boxes ----
    progress_shape.text = progress
    risks_shape.text = risks
    summary_shape.text = summary

    # ---- Force all text to black ----
    for shape in [title_shape, progress_shape, risks_shape, summary_shape]:
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.color.rgb = RGBColor(0, 0, 0)

    # ---- Save the slide ----
    timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
    output_file = f"Dashboard_{timestamp}.pptx"
    prs.save(output_file)
    print(f"✅ Dashboard generated: {output_file}")
    return output_file

# ---- Example usage ----
if __name__ == "__main__":
    exec_summary_text = """**Executive Summary: Program Updates**

**Overall Program Status:** Yellow

**Key Progress Points:**
- The Marketing workstream is nearing completion of all tasks.
- The Engineering team is currently engaged in ongoing integration.
- The Sales pipeline remains strong with no major risks reported.

**Risks and Mitigation Actions:**
- **Marketing:** A timing issue presents a risk, which will be addressed by allocating extra resources.
- **Engineering:** There is a potential delay in deployment, mitigated by adding extra quality assurance resources.
"""
    generate_dashboard_from_exec(exec_summary_text)
