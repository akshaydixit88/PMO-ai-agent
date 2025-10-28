import openai
from storage import get_all_updates
from pptx import Presentation
from pptx.util import Pt

def generate_exec_summary(api_key: str) -> str:
    """
    Fetch updates from the DB, send to OpenAI, and return an executive summary.
    """
    # 1️⃣ Get all updates from the DB
    updates = get_all_updates()  # list of tuples: (workstream, status, progress, risks, mitigation, timestamp)
    
    if not updates:
        return "No updates found for this period."

    # 2️⃣ Format updates nicely for the prompt
    formatted_updates = ""
    for u in updates:
        workstream, status, progress, risks, mitigation, timestamp = u
        formatted_updates += (
            f"Workstream: {workstream}\n"
            f"Status: {status}\n"
            f"Progress: {progress}\n"
            f"Risks: {risks}\n"
            f"Mitigation: {mitigation}\n"
            f"Timestamp: {timestamp}\n\n"
        )

    # 3️⃣ Prepare prompt
    prompt = f"""
You are an executive assistant summarizing program updates.

Here are the workstream updates:

{formatted_updates}

Your job:
- Identify only the most important, organization-wide takeaways. Do not summarize by each workstream.
- Determine the **overall program status** (Green, Yellow, or Red) based on trends and severity of risks.
- Combine progress updates into a concise list of **Key Progress Points** (2–4 bullets maximum).
- Combine risks and mitigations into a concise **Risks and Mitigation Actions** section (2–3 bullets maximum).
- Ignore duplicate or outdated updates for the same workstream; use the latest only.
- Write in polished business language suitable for an executive report (no filler words).
- Output exactly in this format:

**Executive Summary: Program Updates**

**Overall Program Status:** <Green/Yellow/Red>

**Key Progress Points:**
- <Point 1>
- <Point 2>
- <Point 3>

**Risks and Mitigation Actions:**
- <Risk/Mitigation 1>
- <Risk/Mitigation 2>
- <Risk/Mitigation 3>
"""

    # 4️⃣ Call OpenAI API
    client = openai.OpenAI(api_key=api_key)
    response = client.chat.completions.create(
        model="gpt-4o-mini",
        messages=[{"role": "user", "content": prompt}]
    )

    summary = response.choices[0].message.content
    return summary



def create_summary_slide(summary_text: str, output_file: str = "Program_Exec_Summary.pptx"):
    prs = Presentation()
    slide_layout = prs.slide_layouts[1]  # Title + content
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "Program Dashboard Summary"

    textbox = slide.shapes.placeholders[1]  # content placeholder
    textbox.text = summary_text

    # Optional: adjust font size
    for paragraph in textbox.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.size = Pt(14)

    prs.save(output_file)
    print(f"✅ PowerPoint saved as {output_file}")
